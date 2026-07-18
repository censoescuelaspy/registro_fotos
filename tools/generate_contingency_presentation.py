#!/usr/bin/env python3
"""Genera la ficha CIALPA v1.4 editable con croquis y manual de campo."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


FORM_ID = "CIALPA_PM_01"
FORM_VERSION = "1.4"
PAGE_WIDTH_MM = 215.9
PAGE_HEIGHT_MM = 330.2

NAVY = "173F67"
ORANGE = "E94B24"
INK = "172033"
GRAY = "667085"
BORDER = "98A2B3"
GRID_MINOR = "E1E6EB"
GRID_MAJOR = "BBC5CF"
PALE_ORANGE = "FFF7ED"

TYPE_OPTIONS = (
    "□ PG Planta gral.   □ AUL Aula   □ ADM Administración   □ BIB Biblioteca   "
    "□ COM Cocina/Comedor   □ DEP Depósito\n"
    "□ LAB Laboratorio   □ TAL Taller   □ SAN Sanitario   □ CIR Pasillo   "
    "□ EXT Exterior   □ OTR Otro: __________"
)

SYMBOLS = [
    ("MU", "Pared", "wall"),
    ("TA", "Tabique", "partition"),
    ("PT", "Puerta", "door"),
    ("PD", "Puerta dbl.", "double_door"),
    ("VN", "Ventana", "window"),
    ("PI", "Pilar", "pillar"),
    ("CT", "Cota/med.", "dimension"),
    ("ES", "Escalera", "stairs"),
    ("RM", "Rampa", "ramp"),
    ("WC", "Inodoro", "badge"),
    ("LV", "Lavamanos", "badge"),
    ("UR", "Urinario", "badge"),
    ("DU", "Ducha", "badge"),
    ("TC", "Toma", "badge_square"),
    ("IN", "Interruptor", "badge"),
    ("LU", "Luz", "badge"),
    ("VE", "Ventilador", "badge"),
    ("AA", "Aire acond.", "badge_square"),
    ("TE", "Tablero", "badge_square"),
    ("AP", "Pto. agua", "badge"),
    ("DG", "Desagüe", "badge"),
    ("DF", "Daño/falla", "badge"),
    ("FO", "Pto. foto", "badge"),
    ("OT", "Otro", "badge"),
]


def mm(value: float):
    return Cm(value / 10.0)


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_shape_fill(shape, color: str | None) -> None:
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)


def add_rect(slide, x: float, y: float, w: float, h: float, *, line=BORDER, fill=None, width=0.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, mm(x), mm(y), mm(w), mm(h))
    set_shape_fill(shape, fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size=6.0,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.25,
    font="Arial",
):
    shape = slide.shapes.add_textbox(mm(x), mm(y), mm(w), mm(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = mm(margin)
    frame.margin_right = mm(margin)
    frame.margin_top = mm(0.05)
    frame.margin_bottom = mm(0.05)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 0.85
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color=INK, width=0.6, dash=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mm(x1), mm(y1), mm(x2), mm(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_panel(slide, x: float, y: float, w: float, h: float, label: str, value: str, *, value_size=5.0):
    add_rect(slide, x, y, w, h)
    add_text(slide, x + 0.4, y + 0.35, w - 0.8, 2.0, label, size=4.1, color=NAVY, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, x + 0.5, y + 2.8, w - 1.0, h - 3.0, value, size=value_size, bold=True, valign=MSO_ANCHOR.TOP)


def add_box_field(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    count: int,
    *,
    box_w=5.2,
    box_h=5.2,
):
    add_rect(slide, x, y, w, h)
    add_text(slide, x + 0.4, y + 0.35, w - 0.8, 2.0, label, size=4.1, color=NAVY, bold=True, valign=MSO_ANCHOR.TOP)
    start_x = x + 0.7
    start_y = y + 4.0
    for index in range(count):
        add_rect(slide, start_x + index * box_w, start_y, box_w, box_h, line=NAVY, width=0.9)


def add_header(slide, logo: Path | None) -> None:
    add_rect(slide, 9.5, 8.3, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_rect(slide, 204.0, 8.3, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    if logo and logo.exists():
        slide.shapes.add_picture(str(logo), mm(17), mm(5.4), width=mm(11.5))
    else:
        add_text(slide, 15, 6, 18, 7, "CIALPA S.A.", size=7, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        48,
        5.3,
        120,
        8.5,
        "FICHA DE CONTINGENCIA - PLANO MANUAL",
        size=9.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, 177, 5.5, 25, 8, f"FORM {FORM_ID}\nv{FORM_VERSION}", size=5.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def add_identity(slide) -> None:
    x, y, total = 7.0, 16.0, 202.0
    col = total / 30.0
    row0_h, row1_h = 12.8, 12.2
    add_box_field(slide, x, y, col * 9, row0_h, "[E01] CÓDIGO MEC DE ESCUELA", 10)
    add_box_field(slide, x + col * 9, y, col * 9, row0_h, "[E04] CÓDIGO CENSISTA / CÉDULA (SIN PUNTOS)", 10)
    add_box_field(slide, x + col * 18, y, col * 6, row0_h, "[E02] FECHA AAAA-MM-DD", 8, box_w=4.6)
    add_box_field(slide, x + col * 24, y, col * 3, row0_h, "[E08] FORM.", 2)
    add_box_field(slide, x + col * 27, y, col * 3, row0_h, "[E06] HOJA", 2)

    y += row0_h
    add_box_field(slide, x, y, col * 5, row1_h, "[R01] BLOQUE B", 2, box_w=6.0, box_h=5.5)
    add_box_field(slide, x + col * 5, y, col * 5, row1_h, "[R02] PISO P", 2, box_w=6.0, box_h=5.5)
    add_box_field(slide, x + col * 10, y, col * 7, row1_h, "[R03] ESPACIO E", 3, box_w=6.0, box_h=5.5)
    add_box_field(slide, x + col * 17, y, col * 4, row1_h, "[R05] CANTIDAD DE FOTOS", 2, box_w=6.0, box_h=5.5)
    add_panel(
        slide,
        x + col * 21,
        y,
        col * 9,
        row1_h,
        "[E10] ORIENTACIÓN NORTE",
        "□ ↑ Arriba   □ → Derecha   □ ↓ Abajo   □ ← Izquierda   □ Otra",
        value_size=4.6,
    )


def add_record_fields(slide) -> None:
    x, y, total = 7.0, 41.0, 202.0
    add_panel(slide, x, y, total, 6.8, "[R04] TIPO DE ESPACIO - MARQUE UNA OPCIÓN", TYPE_OPTIONS, value_size=3.8)
    y += 6.8
    add_panel(
        slide,
        x,
        y,
        total * 14 / 24,
        8.0,
        "MEDIDAS DIRECTAS EN CENTÍMETROS (cm) - NO CALCULAR ÁREA",
        "[M01] Largo __________ cm    [M02] Ancho __________ cm    [M03] Alto __________ cm",
        value_size=4.7,
    )
    add_panel(
        slide,
        x + total * 14 / 24,
        y,
        total * 10 / 24,
        8.0,
        "[E12] ESTADO DEL CROQUIS - MARQUE UNA OPCIÓN",
        "□ Inicial   □ Rev.   □ Final   □ Parc.",
        value_size=5.0,
    )


def add_instructions(slide) -> None:
    x, y = 7.0, 55.8
    add_rect(slide, x, y, 126, 6.5, line="FDBA74", fill=PALE_ORANGE, width=0.7)
    add_rect(slide, x + 126, y, 76, 6.5, line="FDBA74", fill=PALE_ORANGE, width=0.7)
    add_text(
        slide,
        x + 1,
        y + 0.3,
        124,
        5.9,
        "MEDIDAS: SIEMPRE EN CENTÍMETROS (cm). Use coma decimal solo si es imprescindible expresar milímetros "
        "(ej.: 12,5 cm = 125 mm). No use metros ni calcule áreas.",
        size=4.8,
        color="9A3412",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x + 127,
        y + 0.3,
        74,
        5.9,
        "TINTA: negro o azul oscuro. Rojo opcional solo para resaltar DF; escriba siempre el código DF en tinta oscura.",
        size=4.5,
        color="9A3412",
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_symbol_icon(slide, x: float, y: float, code: str, kind: str) -> None:
    if kind in {"badge", "badge_square"}:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL if kind == "badge" else MSO_SHAPE.RECTANGLE,
            mm(x),
            mm(y),
            mm(3.4),
            mm(3.4),
        )
        set_shape_fill(badge, None)
        badge.line.color.rgb = rgb(INK)
        badge.line.width = Pt(0.55)
        add_text(slide, x, y + 0.05, 3.4, 3.2, code, size=2.5, bold=True, align=PP_ALIGN.CENTER)
        return
    if kind == "wall":
        add_line(slide, x, y + 1.7, x + 4.2, y + 1.7, width=1.4)
    elif kind == "partition":
        add_line(slide, x, y + 1.7, x + 4.2, y + 1.7, width=0.55, dash=True)
    elif kind == "door":
        add_line(slide, x, y + 3.2, x, y + 0.4, width=0.65)
        add_line(slide, x, y + 0.4, x + 2.7, y + 3.2, width=0.65)
        add_line(slide, x + 3.2, y + 3.2, x + 4.2, y + 3.2, width=0.65)
    elif kind == "double_door":
        add_line(slide, x, y + 3.2, x + 2.0, y + 0.5, width=0.65)
        add_line(slide, x + 4.2, y + 3.2, x + 2.2, y + 0.5, width=0.65)
    elif kind == "window":
        add_line(slide, x, y + 1.2, x + 4.2, y + 1.2, width=0.55)
        add_line(slide, x, y + 2.2, x + 4.2, y + 2.2, width=0.55)
        add_line(slide, x + 0.8, y + 0.6, x + 0.8, y + 2.8, width=0.5)
        add_line(slide, x + 3.4, y + 0.6, x + 3.4, y + 2.8, width=0.5)
    elif kind == "pillar":
        add_rect(slide, x + 1.0, y + 0.6, 2.2, 2.2, line=INK, fill=INK, width=0.4)
    elif kind == "dimension":
        add_line(slide, x, y + 1.7, x + 4.2, y + 1.7, width=0.55)
        add_line(slide, x, y + 0.8, x, y + 2.6, width=0.55)
        add_line(slide, x + 4.2, y + 0.8, x + 4.2, y + 2.6, width=0.55)
    elif kind == "stairs":
        for step in range(4):
            sx = x + step * 1.0
            add_line(slide, sx, y + 3.0 - step * 0.65, sx + 1.0, y + 3.0 - step * 0.65, width=0.5)
            if step < 3:
                add_line(slide, sx + 1.0, y + 3.0 - step * 0.65, sx + 1.0, y + 2.35 - step * 0.65, width=0.5)
    elif kind == "ramp":
        add_rect(slide, x, y + 0.5, 4.2, 2.6, line=INK, width=0.5)
        add_line(slide, x + 0.4, y + 2.7, x + 3.7, y + 0.9, width=0.55)


def add_legend(slide) -> None:
    x, y, total = 7.0, 62.3, 202.0
    cell_w, cell_h = total / 8, 5.8
    for index, (code, label, kind) in enumerate(SYMBOLS):
        row, col = divmod(index, 8)
        cell_x, cell_y = x + col * cell_w, y + row * cell_h
        add_rect(slide, cell_x, cell_y, cell_w, cell_h, line="D0D5DD", width=0.45)
        add_symbol_icon(slide, cell_x + 1.2, cell_y + 1.2, code, kind)
        add_text(slide, cell_x + 5.9, cell_y + 0.6, cell_w - 6.4, 4.6, f"{code} {label}", size=4.0, bold=True)


def add_grid(slide) -> None:
    x, header_y, grid_y = 7.0, 79.7, 84.2
    add_text(slide, x + 0.8, header_y, 55, 4.5, "PLANO / CROQUIS", size=5.7, color=NAVY, bold=True)
    add_text(
        slide,
        x + 55,
        header_y,
        100,
        4.5,
        "ESCALA: □ 1 cuadro=50 cm   □ 1 cuadro=100 cm   □ sin escala",
        size=5.0,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, x + 155, header_y, 46, 4.5, "NORTE: ↑ / dibuje N", size=5.0, bold=True, align=PP_ALIGN.RIGHT)

    grid_x, grid_w, grid_h = 8.0, 200.0, 210.0
    for step in range(41):
        gx = grid_x + step * 5.0
        major = step % 5 == 0
        add_line(slide, gx, grid_y, gx, grid_y + grid_h, color=GRID_MAJOR if major else GRID_MINOR, width=0.7 if major else 0.25)
    for step in range(43):
        gy = grid_y + step * 5.0
        major = step % 5 == 0
        add_line(slide, grid_x, gy, grid_x + grid_w, gy, color=GRID_MAJOR if major else GRID_MINOR, width=0.7 if major else 0.25)
    add_rect(slide, grid_x, grid_y, grid_w, grid_h, line=GRAY, width=1.0)


def add_summary(slide) -> None:
    x, y, total = 7.0, 294.2, 202.0
    labels = [
        "[C01]PT___", "[C02]VN___", "[C03]PI___", "[C04]WC___", "[C05]LV___", "[C06]UR___", "[C07]DU___",
        "[C08]TC___", "[C09]IN___", "[C10]LU___", "[C11]VE___", "[C12]AA___", "[C13]TE___", "[C14]DF___",
    ]
    cell_w = total / len(labels)
    for index, label in enumerate(labels):
        add_rect(slide, x + index * cell_w, y, cell_w, 5.0, line=BORDER, width=0.5)
        add_text(slide, x + index * cell_w, y, cell_w, 5.0, label, size=3.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    y += 5.0
    add_rect(slide, x, y, total, 9.0, line=BORDER, width=0.6)
    add_text(
        slide,
        x + 0.8,
        y + 0.4,
        total - 1.6,
        2.2,
        "[O01] OBSERVACIONES / DAÑOS Y FALLAS (use DF01, DF02... y vincule F01, F02...):",
        size=4.5,
        color=NAVY,
        bold=True,
    )
    add_line(slide, x + 2, y + 4.6, x + 112, y + 4.6, color=BORDER, width=0.45)
    add_line(slide, x + 2, y + 7.1, x + 112, y + 7.1, color=BORDER, width=0.45)
    add_text(
        slide,
        x + 116,
        y + 0.5,
        84,
        3.0,
        "REFERENCIA DE FOTOS EN EL CROQUIS",
        size=4.2,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x + 116,
        y + 3.0,
        84,
        5.3,
        "Marque FO01, FO02... junto al elemento. En la app seleccione su código y número, "
        "por ejemplo PT01 o DF02. La identificación completa se imprime al pie de la foto.",
        size=3.9,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_footer(slide) -> None:
    add_rect(slide, 9.5, 314.2, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_rect(slide, 204.0, 314.2, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_text(
        slide,
        37,
        311.5,
        142,
        7.2,
        f"{FORM_ID} v{FORM_VERSION} | Hoja 1 de 2 | ID FOTO: ESCUELA-B##-P##-E###-H##-ELEMENTO##-FT## | "
        "[E07] ESTADO: □ COMPLETO  □ PARCIAL | Fotografíe la hoja completa, plana, enfocada y sin sombras.",
        size=4.15,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_manual_header(slide, logo: Path | None) -> None:
    add_rect(slide, 9.5, 8.3, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_rect(slide, 204.0, 8.3, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    if logo and logo.exists():
        slide.shapes.add_picture(str(logo), mm(16), mm(5.3), width=mm(12.0))
    add_text(
        slide,
        40,
        5.0,
        136,
        8.0,
        "MANUAL RÁPIDO DEL CENSISTA",
        size=10.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        173,
        5.4,
        30,
        7.5,
        f"FORM {FORM_ID}\nv{FORM_VERSION} · HOJA 2/2",
        size=5.3,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_manual_title(slide, x: float, y: float, w: float, number: str, title: str, subtitle: str = "") -> None:
    add_rect(slide, x, y, w, 7.0, line=NAVY, fill=NAVY, width=0.5)
    add_text(slide, x + 1.0, y + 0.4, 8.0, 6.0, number, size=8.0, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 9.0, y + 0.3, w - 10.0, 3.2, title, size=6.2, color="FFFFFF", bold=True)
    if subtitle:
        add_text(slide, x + 9.0, y + 3.1, w - 10.0, 3.0, subtitle, size=3.8, color="D9E7F2")


def add_manual_step(slide, x: float, y: float, w: float, number: str, title: str, body: str) -> None:
    add_rect(slide, x, y, w, 29.0, line="CBD5E1", fill="F8FAFC", width=0.55)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, mm(x + 2.0), mm(y + 2.2), mm(7.0), mm(7.0))
    set_shape_fill(badge, ORANGE)
    badge.line.color.rgb = rgb(ORANGE)
    add_text(slide, x + 2.0, y + 2.3, 7.0, 6.6, number, size=6.0, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 10.2, y + 2.0, w - 12.0, 6.0, title, size=5.8, color=NAVY, bold=True)
    add_text(slide, x + 2.2, y + 10.0, w - 4.4, 16.5, body, size=5.15, color=INK, valign=MSO_ANCHOR.TOP)


def add_manual_slide(slide, logo: Path | None) -> None:
    add_manual_header(slide, logo)
    x, width = 7.0, 202.0

    add_manual_title(
        slide,
        x,
        16.0,
        width,
        "A",
        "SECUENCIA DE TRABAJO",
        "Complete el papel y la app con los mismos números. No improvise códigos.",
    )
    step_y = 24.0
    gap = 2.0
    step_w = (width - gap * 2) / 3
    steps = [
        ("1", "UBÍQUESE", "Ingrese a CIALPA Fotos, permita ubicación y abra la escuela asignada desde el mapa o la lista."),
        ("2", "IDENTIFIQUE LA HOJA", "Escriba código MEC, cédula, fecha, formulario, hoja, bloque, piso y espacio antes de dibujar."),
        ("3", "DIBUJE Y NUMERE", "Trace el croquis. Use los símbolos oficiales y numere cada tipo por espacio: PT01, PT02; VN01; DF01..."),
        ("4", "REPITA LOS DATOS", "En la app elija exactamente la misma escuela, formulario, hoja, bloque, piso, espacio y tipo de espacio."),
        ("5", "FOTOGRAFÍE", "Seleccione tipo y número de elemento. Pulse Foto del espacio. La app añade el código completo al pie sin tapar la imagen."),
        ("6", "CIERRE Y SINCRONICE", "Tome una foto completa de la hoja, revise nitidez y códigos, finalice y confirme que Pendientes quede en cero."),
    ]
    for index, (number, title, body) in enumerate(steps):
        row, col = divmod(index, 3)
        add_manual_step(slide, x + col * (step_w + gap), step_y + row * 31.0, step_w, number, title, body)

    add_manual_title(slide, x, 87.0, width, "B", "CÓDIGO ÚNICO DE LA FOTO", "Debe coincidir en papel, imagen, Drive y Sheets.")
    add_rect(slide, x, 94.0, width, 22.0, line="9BC8B9", fill="E8F6F1", width=0.7)
    add_text(
        slide,
        x + 3.0,
        96.0,
        width - 6.0,
        7.0,
        "11007-B01-P00-E001-H01-PT01-FT01",
        size=10.2,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        font="Consolas",
    )
    add_text(
        slide,
        x + 3.0,
        104.0,
        width - 6.0,
        9.5,
        "ESCUELA 11007  |  BLOQUE 01  |  PISO 00 (planta baja)  |  ESPACIO 001  |  HOJA 01  |  PUERTA 01  |  FOTO 01",
        size=4.7,
        color="164E3F",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_manual_title(slide, x, 119.0, 98.5, "C", "CÓMO NUMERAR", "Reinicie por espacio y por tipo.")
    add_manual_title(slide, x + 101.5, 119.0, 100.5, "D", "FOTOS QUE SÍ SIRVEN", "Primero contexto; luego detalle.")

    numbering = [
        ("B01", "Primer bloque de la escuela"),
        ("P00", "Planta baja; P01 primer piso"),
        ("E001", "Primer espacio del bloque y piso"),
        ("PT01", "Primera puerta de ese espacio"),
        ("DF01", "Primer daño o falla señalado"),
        ("FO01", "Punto desde donde se tomó la foto"),
    ]
    panel_y = 126.0
    add_rect(slide, x, panel_y, 98.5, 58.0, line="CBD5E1", fill="FFFFFF", width=0.55)
    for index, (code, meaning) in enumerate(numbering):
        row_y = panel_y + 2.0 + index * 9.0
        add_rect(slide, x + 2.0, row_y, 15.0, 6.8, line="9FB3C5", fill="EDF4F9", width=0.45)
        add_text(slide, x + 2.0, row_y, 15.0, 6.8, code, size=5.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
        add_text(slide, x + 19.0, row_y, 76.0, 6.8, meaning, size=4.9, color=INK)

    photo_rules = [
        ("1. CONTEXTO", "Muestre el elemento completo y su relación con el espacio."),
        ("2. DETALLE", "Para daños, agregue una toma cercana después de la vista general."),
        ("3. HOJA", "Cuatro esquinas visibles, cámara paralela, sin sombra y texto enfocado."),
        ("4. EVITE", "Dedos, contraluz, zoom digital, fotos repetidas o números distintos."),
    ]
    add_rect(slide, x + 101.5, panel_y, 100.5, 58.0, line="CBD5E1", fill="FFFFFF", width=0.55)
    for index, (title, body) in enumerate(photo_rules):
        row_y = panel_y + 2.0 + index * 13.5
        add_text(slide, x + 104.0, row_y, 22.0, 5.5, title, size=4.4, color=ORANGE, bold=True)
        add_text(slide, x + 126.0, row_y, 73.0, 10.5, body, size=4.85, color=INK, valign=MSO_ANCHOR.TOP)

    add_manual_title(slide, x, 187.0, width, "E", "SI NO HAY INTERNET", "La cámara y los borradores siguen funcionando.")
    add_rect(slide, x, 194.0, width, 24.0, line="F0B86E", fill=PALE_ORANGE, width=0.7)
    add_text(
        slide,
        x + 3.0,
        196.0,
        width - 6.0,
        20.0,
        "1. Pulse Guardar borrador o Finalizar y guardar en cola.  2. No cierre sesión, no desinstale la app y no borre datos del navegador.  "
        "3. Cuando vuelva la señal, abra Pendientes y pulse Sincronizar ahora.  4. Espere hasta ver 0 operaciones en cola.",
        size=5.0,
        color="8A4510",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_manual_title(slide, x, 221.0, width, "F", "CONTROL FINAL ANTES DE RETIRARSE", "Revise cada casilla; una foto sin relación no puede recuperarse después.")
    checks = [
        "□ Código MEC correcto",
        "□ Cédula sin puntos",
        "□ B/P/E/H iguales en papel y app",
        "□ Elementos numerados",
        "□ Medidas solo en cm",
        "□ Daños/fallas marcados DF",
        "□ Hoja completa fotografiada",
        "□ Pendientes en cero o guardados",
    ]
    check_y = 228.0
    for index, check in enumerate(checks):
        row, col = divmod(index, 2)
        cell_x = x + col * 101.0
        cell_y = check_y + row * 10.0
        add_rect(slide, cell_x, cell_y, 101.0, 10.0, line="D0D5DD", fill="FFFFFF", width=0.45)
        add_text(slide, cell_x + 2.0, cell_y + 0.8, 97.0, 8.2, check, size=4.8, color=INK, bold=True)

    add_rect(slide, x, 270.0, width, 30.0, line=NAVY, fill="F2F7FB", width=0.7)
    add_text(slide, x + 2.0, 272.0, 62.0, 5.0, "REGLAS QUE NO CAMBIAN", size=5.2, color=NAVY, bold=True)
    add_text(
        slide,
        x + 2.0,
        278.0,
        96.0,
        19.0,
        "• Todas las medidas se anotan en centímetros (cm).\n"
        "• Use coma decimal solo para milímetros: 12,5 cm.\n"
        "• Tinta negra o azul oscura; rojo solo resalta DF.\n"
        "• No borre: tache una vez y escriba la corrección al lado.",
        size=4.9,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(slide, x + 105.0, 272.0, 94.0, 5.0, "ACCESO A LA APP", size=5.2, color=NAVY, bold=True)
    add_text(
        slide,
        x + 105.0,
        278.0,
        94.0,
        17.0,
        "censoescuelaspy.github.io/registro_fotos/\n\n"
        "Si una escuela no aparece o el usuario no puede ingresar, conserve la ficha y comuníquese con el administrador antes de crear códigos alternativos.",
        size=4.75,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )

    add_rect(slide, 9.5, 314.2, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_rect(slide, 204.0, 314.2, 2.4, 2.4, line="000000", fill="000000", width=0.3)
    add_text(
        slide,
        32,
        307.0,
        152,
        11.0,
        f"{FORM_ID} v{FORM_VERSION} | Hoja 2 de 2 | Conserve esta hoja junto a las fichas en blanco. "
        "La identificación impresa por la app no reemplaza la numeración correcta del croquis.",
        size=4.6,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def generate(output: Path, logo: Path | None) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = mm(PAGE_WIDTH_MM)
    presentation.slide_height = mm(PAGE_HEIGHT_MM)
    presentation.core_properties.title = "Ficha de contingencia y manual de campo CIALPA v1.4"
    presentation.core_properties.subject = "Versión editable para Google Slides"
    presentation.core_properties.author = "CIALPA"
    presentation.core_properties.keywords = "CIALPA, plano manual, editable, Google Slides, OCR"

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb("FFFFFF")

    add_header(slide, logo)
    add_identity(slide)
    add_record_fields(slide)
    add_instructions(slide)
    add_legend(slide)
    add_grid(slide)
    add_summary(slide)
    add_footer(slide)

    manual_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    manual_background = manual_slide.background.fill
    manual_background.solid()
    manual_background.fore_color.rgb = rgb("FFFFFF")
    add_manual_slide(manual_slide, logo)
    presentation.save(output)


def parse_args() -> argparse.Namespace:
    project_root = Path(__file__).resolve().parents[1]
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=project_root
        / "docs"
        / f"FICHA_CONTINGENCIA_PLANO_MANUAL_CIALPA_v{FORM_VERSION}_EDITABLE.pptx",
    )
    parser.add_argument(
        "--logo",
        type=Path,
        default=project_root / "assets" / "img" / "logo.png",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    logo = args.logo if args.logo.exists() else None
    generate(args.output.resolve(), logo)
    print(f"PPTX: {args.output.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
